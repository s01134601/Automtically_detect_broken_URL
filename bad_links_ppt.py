import os
import requests
from pptx import Presentation
import pandas as pd
from paths import source_dir, target_dir_2


def gen_doc_list(source) -> list:
    """
    returns relevant files from source directory
    :param source: path to directory with relevant files
    :return: list of relevant files
    """
    doc_list = []
    for root, dirnames, filenames in os.walk(source):
        for filename in filenames:
            if filename.endswith(('.ppt', '.pptx')):
                doc_list.append(os.path.join(root, filename))
    return doc_list


def get_url_status(doc_url_df):
    """
    :param doc_url_df: a dataframe that has the following columns: url, path, page number
    :return: an additional column to the input dataframe
    """
    stat_codes = []
    for index, row in doc_url_df.iterrows():
        try:
            response = requests.get(row.url, verify=False, timeout=10.0)
            stat_codes.append(response.status_code)
        except requests.exceptions.Timeout:
            stat_codes.append(504)
        except requests.exceptions.ConnectionError:
            stat_codes.append(-1)
    doc_url_df['status'] = stat_codes
    return doc_url_df


if '__main__' == __name__:

    # find doc and docx files
    requests.packages.urllib3.disable_warnings()
    doc_list = gen_doc_list(source_dir)
    print(doc_list)

url_list = []
page_n = []
ppt_name = []
for ppt in doc_list:
    ppt_df = pd.DataFrame()
    """
    The following just goes through each slide, picks out text frames, and try to run the hyperlink
    associated with the texts. It outputs all the hyperlinks associated with texts. 
    """
    prs = Presentation(ppt)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    address = run.hyperlink.address
                    if not address is None:
                        url_list.append(address)
                        page_n.append(prs.slides.index(slide))
                        ppt_name.append(ppt)
ppt_df['path'] = ppt_name
ppt_df['page_n'] = page_n
ppt_df['url'] = url_list


requests.packages.urllib3.disable_warnings()
ppt_df = get_url_status(ppt_df)
# Pick out only bad links (ones that don't have 200 as status code)
ppt_urls = ppt_df.loc[ppt_df['status']!= 200]
ppt_urls.to_csv(target_dir_2)